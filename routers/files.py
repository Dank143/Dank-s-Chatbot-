import io

from fastapi import APIRouter, File, HTTPException, UploadFile

router = APIRouter(prefix="/api")

_MAX_DOC_CHARS = 120_000


@router.post("/extract-text")
async def extract_text(file: UploadFile = File(...)):
    """Extract plain text from an uploaded txt/pdf/pptx/docx/xlsx (capped length)."""
    content = await file.read()
    filename = file.filename or "file"
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    try:
        if ext == "txt":
            text = content.decode("utf-8", errors="replace")

        elif ext == "pdf":
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(content))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)

        elif ext == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            text = "\n".join(parts)

        elif ext == "docx":
            from docx import Document
            doc = Document(io.BytesIO(content))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif ext == "xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                rows.append(f"[Sheet: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join("" if cell is None else str(cell) for cell in row)
                    if row_text.strip():
                        rows.append(row_text)
            text = "\n".join(rows)

        else:
            raise HTTPException(400, f"Unsupported file type: .{ext}")

    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(422, f"Could not read file: {exc}")

    text = text.strip()
    if len(text) > _MAX_DOC_CHARS:
        text = text[:_MAX_DOC_CHARS] + f"\n\n[Truncated — file exceeds {_MAX_DOC_CHARS:,} characters]"

    return {"text": text, "filename": filename}
